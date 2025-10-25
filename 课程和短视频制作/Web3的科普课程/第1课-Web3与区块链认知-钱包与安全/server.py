#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
轻量后端：
1. 托管当前目录静态文件（含 HTML）
2. GET /download-ppt  → 返回与 HTML 同内容与备注的 .pptx
"""
import os, io, re
from http.server import HTTPServer, SimpleHTTPRequestHandler
from pptx import Presentation
from pptx.util import Inches, Pt

ROOT = os.path.dirname(os.path.abspath(__file__))
MD_FILE = os.path.join(ROOT, '第1课-逐字稿.md')

# 读取逐字稿，按页拆标题+口播
def load_notes():
    with open(MD_FILE, encoding='utf-8') as f:
        md = f.read()
    blocks = re.findall(r'【(P\d+).*?】[\s\S]*?逐字稿（口播）：([\s\S]*?)(?=\n\n)', md)
    return {p:n.replace('\n',' ').strip() for p,n in blocks}

# 生成 PPT
def build_ppt():
    prs = Presentation()
    notes = load_notes()
    # 读取 HTML 里的标题与 bullet
    with open(os.path.join(ROOT, '第1课-PPT.html'), encoding='utf-8') as f:
        html = f.read()
    slides_data = re.findall(r'<h2>(P\d+.*?)<\/h2>[\s\S]*?<div class="content">([\s\S]*?)<\/div>', html)
    for title, content in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title & Content
        slide.shapes.title.text = title
        body = slide.shapes.placeholders[1].text_frame
        body.clear()
        for line in content.split('<br>'):
            p = body.add_paragraph()
            p.text = line.strip()
            p.font.size = Pt(18)
        # 加备注
        if title[:2] in notes:
            slide.notes_slide.notes_text_frame.text = notes[title[:2]]
    bio = io.BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio

class Handler(SimpleHTTPRequestHandler):
    def do_GET(self):
        if self.path == '/download-ppt':
            bio = build_ppt()
            self.send_response(200)
            self.send_header('Content-Type', 'application/vnd.openxmlformats-officedocument.presentationml.presentation')
            self.send_header('Content-Disposition', 'attachment; filename="第1课-Web3钱包与安全.pptx"')
            self.end_headers()
            self.wfile.write(bio.read())
        else:
            super().do_GET()

if __name__ == '__main__':
    os.chdir(ROOT)
    srv = HTTPServer(('0.0.0.0', 8080), Handler)
    print('Serving at http://0.0.0.0:8080  …')
    srv.serve_forever()